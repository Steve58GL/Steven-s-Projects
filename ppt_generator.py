import requests
import pandas as pd
from datetime import datetime
import win32com.client
import os
import json
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches



datetime_now = datetime.now()
full_list_url = 'https://finance.yahoo.com/currencies'
full_list_page = requests.get(full_list_url)
df = pd.read_html(full_list_page.text)[0].drop_duplicates()
df['pct_change'] = df['% Change'].str.slice(stop=-1).astype(float)

top_df = df.sort_values(['pct_change'], ascending=False).reset_index(drop=True)[:5]
top_df = top_df[['Name', 'Last Price', 'Change', '% Change']]

bottom_df = df.sort_values(['pct_change'], ascending=True).reset_index(drop=True)[:5]
bottom_df = bottom_df[['Name', 'Last Price', 'Change', '% Change']]

top_name = top_df['Name'][0].replace('/', '')
bottom_name = bottom_df['Name'][0].replace('/', '')

for idx in range(2):

    name = [top_name, bottom_name][idx]
    file_path = ['top.png', 'bottom.png']

    url = 'https://query1.finance.yahoo.com/v8/finance/chart/' + name + '=X?region=US&lang=en-US&includePrePost=false&interval=30m&useYfid=true&range=1mo&corsDomain=finance.yahoo.com&.tsrc=finance'
    header = {
        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_11_5) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/50.0.2661.102 Safari/537.36'}
    page = requests.get(url, headers=header)
    temp_json = json.loads(page.text)
    price_list = temp_json['chart']['result'][0]['indicators']['quote'][0]['close']
    price_list = [price for price in price_list if price is not None]
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.plot(price_list, color='#43B7A4')
    ax.set_xticks([])
    ax.tick_params(axis='y', colors='#43B7A4', labelsize=20)
    for axis in ['top', 'bottom', 'left', 'right']:
        ax.spines[axis].set_color('#43B7A4')
        ax.spines[axis].set_linewidth(4)
    plt.savefig(file_path[idx], transparent=True)

# Open the PPT
currencies_ppt = Presentation('Currencies.pptx')

# Select the slide to be edited
slide = currencies_ppt.slides[0]

# Remove the old figures
shapes = slide.shapes
for shape in shapes:
    # print(shape.shape_type)
    if shape.shape_type == 13:  # 13 = PICTURE
        shapes.element.remove(shape.element)

# Add the new figures
top_img_path = 'top.png'
bottom_img_path = 'bottom.png'
top_pic = slide.shapes.add_picture(top_img_path, Inches(0.40), Inches(4.85), width=Inches(5.30))
bottom_pic = slide.shapes.add_picture(bottom_img_path, Inches(5.25), Inches(4.85), width=Inches(5.30))

# Send the figures to the back
ref_element = slide.shapes[0]._element
ref_element.addprevious(top_pic._element)
ref_element.addprevious(bottom_pic._element)

# Separate text box and table
shapes = slide.shapes
text_box_list = []
auto_shape_list = []
table_list = []
for shape_idx in range(len(shapes)):
    shape = shapes[shape_idx]
    if shape.shape_type == 17:  # TEXT_BOX
        text_box_list.append(shape_idx)
    if shape.shape_type == 1:  # AUTO_SHAPE
        auto_shape_list.append(shape_idx)
    if shape.shape_type == 19:  # TABLE
        table_list.append(shape_idx)

# Last update date shape index
last_update_date_textbox_height = max([shapes[shape_idx].height for shape_idx in text_box_list])
last_update_date_idx = \
    [shape_idx for shape_idx in text_box_list if shapes[shape_idx].height == last_update_date_textbox_height][0]

# Top 5 figure label shape index
top_label_left = min([shapes[shape_idx].left for shape_idx in auto_shape_list])
top_label_idx = [shape_idx for shape_idx in auto_shape_list if shapes[shape_idx].left == top_label_left][0]
auto_shape_list.remove(top_label_idx)

# Bottom 5 figure label shape index
bottom_label_idx = auto_shape_list[0]

# Top 5 table shape index
top_table_left = min([shapes[shape_idx].left for shape_idx in table_list])
top_table_idx = [shape_idx for shape_idx in table_list if shapes[shape_idx].left == top_table_left][0]
table_list.remove(top_table_idx)

# Bottom 5 table shape index
bottom_table_idx = table_list[0]

# Update last update date
paragraph = shapes[last_update_date_idx].text_frame.paragraphs[0]
paragraph.runs[4].text = datetime_now.strftime("%#d %b %Y %H:%M")

# Update top 5 figure label
paragraph = shapes[top_label_idx].text_frame.paragraphs[0]
paragraph.runs[0].text = top_df['Name'][0].replace('/', ' / ')

# Update bottom 5 figure label
paragraph = shapes[bottom_label_idx].text_frame.paragraphs[0]
paragraph.runs[0].text = bottom_df['Name'][0].replace('/', ' / ')

# Update top table
top_table = shapes[top_table_idx].table
for i in range(5):
    for j in range(4):
        cell = top_table.cell(i + 1, j)
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        run.text = str(top_df.iloc[i, j])

# Update bottom table
bottom_table = shapes[bottom_table_idx].table
for i in range(5):
    for j in range(4):
        cell = bottom_table.cell(i + 1, j)
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        run.text = str(bottom_df.iloc[i, j])

# Save the PPT
currencies_ppt.save('New_Currencies.pptx')

# Open the PPT
ppt_file_path = os.getcwd() + '\\New_Currencies.pptx'
powerpoint = win32com.client.Dispatch('Powerpoint.Application')
deck = powerpoint.Presentations.Open(ppt_file_path)

# Save the PNG
img_file_path = os.getcwd() + '\\Currencies.png'
powerpoint.ActivePresentation.Slides[0].Export(img_file_path, 'PNG')

# Save the PDF
pdf_file_path = os.getcwd() + '\\Currencies.pdf'
deck.SaveAs(pdf_file_path, 32)

# Quit the PPT
deck.Close()
powerpoint.Quit()
os.system('taskkill /F /IM POWERPNT.EXE')
